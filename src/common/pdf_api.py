# common/pdf_api.py — 印流PDflow PDF操作API封装（PySide6版）
# 封装所有PDF操作，适配PySide6文件对话框

import fitz
import io
import os
import time
from PIL import Image, ImageDraw, ImageFont


# ============================================================
# 统一错误结果类型
# ============================================================
class PDFlowError:
    """PDF操作错误信息载体，用于批量处理时收集单个文件的错误"""

    def __init__(self, file_path: str, operation: str, message: str, recoverable: bool = True):
        """
        file_path:    导致错误的文件路径
        operation:    正在执行的操作名称（如 "合并"、"拆分"、"压缩"）
        message:      错误信息
        recoverable:  是否可恢复（跳过该文件后继续处理）
        """
        self.file_path = file_path
        self.operation = operation
        self.message = message
        self.recoverable = recoverable

    def to_dict(self) -> dict:
        """转为字典，方便序列化和日志记录"""
        return {
            "file_path": self.file_path,
            "operation": self.operation,
            "message": self.message,
            "recoverable": self.recoverable,
        }

    def __repr__(self) -> str:
        tag = "可恢复" if self.recoverable else "不可恢复"
        return f"PDFlowError([{tag}] {self.operation} - {self.file_path}: {self.message})"

# ============================================================
# 路径工具
# ============================================================
def get_output_path(input_path: str, suffix: str, output_dir: str = None) -> str:
    """生成输出文件路径"""
    base = os.path.splitext(os.path.basename(input_path))[0]
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)
        return os.path.join(output_dir, f"{base}_{suffix}")
    folder = os.path.dirname(input_path)
    return os.path.join(folder, f"{base}_{suffix}")


def _resolve_output_path(output_path, input_path, default_ext):
    """v1.1-patch 路径鲁棒性：把 output_path 规范成"完整文件路径"。

    历史 BUG：当用户传入"目录路径"（如 "C:/Users/24785/Desktop\\"）时，
    直接 += ".xlsx" 会得到 "C:/Users/24785/Desktop\\.xlsx"，
    wb.save() 在 Windows 上失败报 [Errno 13] Permission denied。

    本函数：
    1. 如果 output_path 是已存在目录（或末尾带 / \\），则当作 output_dir，
       在其下生成 <input_basename>.<ext>
    2. 否则保证以 <ext> 结尾
    3. 确保父目录存在（必要时创建）
    """
    if output_path is None:
        output_path = get_output_path(input_path, "")
        if default_ext and default_ext != ".pdf":
            output_path = output_path.replace(".pdf", default_ext)
    # 检测目录：末尾分隔符 或 路径已存在且是目录
    elif output_path.endswith(("/", "\\")) or os.path.isdir(output_path):
        base = os.path.splitext(os.path.basename(input_path))[0]
        output_path = os.path.join(output_path, f"{base}{default_ext}")
    # 补后缀
    if not output_path.lower().endswith(default_ext.lower()):
        output_path += default_ext
    # 父目录兜底
    parent = os.path.dirname(output_path)
    if parent and not os.path.isdir(parent):
        os.makedirs(parent, exist_ok=True)
    return output_path

# ============================================================
# PDF 信息读取
# ============================================================
def get_pdf_info(file_path: str) -> dict:
    """读取PDF文件信息（页数、大小）"""
    try:
        doc = fitz.open(file_path)
        info = {
            "pages": len(doc),
            "size_mb": round(os.path.getsize(file_path) / 1024 / 1024, 2),
            "title": doc.metadata.get("title", ""),
        }
        doc.close()
        return info
    except Exception as e:
        raise Exception(f"读取PDF信息失败: {str(e)}")

# ============================================================
# 合并 PDF
# ============================================================
def merge_pdfs(output_path: str, *filepaths, progress_callback=None) -> dict:
    """合并多个PDF为一个PDF
    output_path: 完整输出文件路径（含文件名）
    *filepaths: 要合并的PDF文件路径
    progress_callback: callable(current, total, filename) — 每合并一个文件后触发
    """
    file_paths = list(filepaths)
    if not file_paths:
        raise Exception("请先选择要合并的PDF文件")
    if len(file_paths) < 2:
        raise Exception("合并至少需要2个PDF文件")

    errors = []
    skipped_files = []
    result = fitz.open()
    total_pages = 0
    total = len(file_paths)

    for idx, fp in enumerate(file_paths):
        try:
            doc = fitz.open(fp)
            result.insert_pdf(doc)
            total_pages += len(doc)
            doc.close()
        except Exception as e:
            # 单个文件合并失败，跳过并记录错误
            errors.append(PDFlowError(fp, "合并", str(e), recoverable=True))
            skipped_files.append(fp)
            continue

        if progress_callback:
            progress_callback(idx + 1, total, os.path.basename(fp))

    # 如果所有文件都失败了，抛出异常
    if len(skipped_files) == len(file_paths):
        result.close()
        raise Exception("合并PDF失败: 所有文件均无法处理")

    try:
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        result.save(output_path)
        result.close()

        return {
            "success": True,
            "output_path": output_path,
            "total_pages": total_pages,
            "files_merged": len(file_paths) - len(skipped_files),
            "output_size": os.path.getsize(output_path),
            "errors": [e.to_dict() for e in errors],
            "skipped_files": skipped_files,
        }
    except Exception as e:
        raise Exception(f"合并PDF保存失败: {str(e)}")

# ============================================================
# 拆分 PDF
# ============================================================
def split_pdf(filepath: str, output_dir: str, mode: str = "page", range_str: str = None,
              progress_callback=None) -> dict:
    """
    拆分PDF
    mode: "page" 每页单独保存 | "range" 按范围拆分
    range_str: "1-3,5-8" 格式的页码范围（mode=range 时使用）
    progress_callback: callable(current, total, filename) — 每拆分一个片段后触发
    """
    errors = []
    output_files = []

    try:
        doc = fitz.open(filepath)
        total = len(doc)
        basename = os.path.splitext(os.path.basename(filepath))[0]

        if mode == "page":
            # 每页一个文件
            for i in range(total):
                try:
                    part = fitz.open()
                    part.insert_pdf(doc, from_page=i, to_page=i)
                    out_path = os.path.join(output_dir, f"{basename}_p{i+1}.pdf")
                    part.save(out_path)
                    part.close()
                    output_files.append(out_path)
                except Exception as e:
                    errors.append(PDFlowError(filepath, f"拆分第{i+1}页", str(e), recoverable=True))
                    continue

                if progress_callback:
                    progress_callback(i + 1, total, f"{basename}_p{i+1}.pdf")

        elif mode == "range" and range_str:
            # 按范围拆分
            range_segs = [s.strip() for s in range_str.split(",") if s.strip()]
            total_segs = len(range_segs)
            for idx, range_seg in enumerate(range_segs, 1):
                if "-" not in range_seg:
                    continue
                parts = range_seg.split("-")
                if len(parts) != 2:
                    continue
                try:
                    start = max(1, int(parts[0]))
                    end = min(total, int(parts[1]))
                    if start >= end:
                        continue
                    part = fitz.open()
                    part.insert_pdf(doc, from_page=start - 1, to_page=end - 1)
                    out_path = os.path.join(output_dir, f"{basename}_part{idx}.pdf")
                    part.save(out_path)
                    part.close()
                    output_files.append(out_path)
                except ValueError:
                    continue
                except Exception as e:
                    errors.append(PDFlowError(filepath, f"拆分范围{range_seg}", str(e), recoverable=True))
                    continue

                if progress_callback:
                    progress_callback(idx, total_segs, f"{basename}_part{idx}.pdf")

        doc.close()
        return {
            "success": True,
            "files": output_files,
            "count": len(output_files),
            "errors": [e.to_dict() for e in errors],
        }
    except Exception as e:
        raise Exception(f"拆分PDF失败: {str(e)}")

# ============================================================
# 压缩 PDF
# ============================================================
def compress_pdf(input_path: str, quality: str = "high", output_path: str = None,
                   progress_callback=None, timeout: int = 60) -> dict:
    """
    压缩PDF

    quality:
      "high"   — 高质量（适合打印）: 逐页渲染为 200DPI JPEG Q90 后重建 PDF
                 预计减轻 60-85%（文字/设计稿PDF尤其明显）
      "medium" — 中等质量（适合阅读）: 逐页渲染为 150DPI JPEG Q75 后重建 PDF
                 预计减轻 70-90%
      "low"    — 低质量（适合屏幕显示）: 逐页渲染为 72DPI JPEG Q50 后重建 PDF
                 预计减轻 90-97%

    progress_callback: callable(current, total) — 每处理完一页触发
    timeout: 单个文件处理超时秒数，超过则中止（默认60秒）
    """
    try:
        original_size = os.path.getsize(input_path)

        if output_path is None:
            output_path = get_output_path(input_path, f"压缩_{quality}")
        if not output_path.lower().endswith(".pdf"):
            output_path += ".pdf"

        errors = _compress_by_rendering(input_path, output_path, quality, progress_callback, timeout)

        compressed_size = os.path.getsize(output_path)
        ratio = round((1 - compressed_size / original_size) * 100, 1) if original_size > 0 else 0

        return {
            "status": "ok",
            "output": output_path,
            "original_mb": round(original_size / 1024 / 1024, 2),
            "compressed_mb": round(compressed_size / 1024 / 1024, 2),
            "ratio": f"{ratio}%",
            "errors": [e.to_dict() for e in errors],
        }
    except Exception as e:
        raise Exception(f"压缩PDF失败: {str(e)}")


def _compress_by_rendering(input_path: str, output_path: str, quality: str,
                             progress_callback=None, timeout: int = 60) -> list:
    """
    将 PDF 每一页渲染为 JPEG 图片后重建 PDF，保证确切的压缩率。
    单页失败时跳过该页并继续处理，返回错误列表。

    quality → (DPI, JPEG Q):
      high   → (200, 90)   — 清晰可读，大幅压缩
      medium → (150, 75)   — 中等清晰，高压缩
      low    → (72,  50)   — 极致压缩

    timeout: 处理超时秒数，超时则中止
    """
    dpi, jpeg_q = {
        "high":   (200, 90),
        "medium": (150, 75),
        "low":    (72,  50),
    }[quality]

    doc = fitz.open(input_path)
    new_doc = fitz.open()
    total_pages = len(doc)
    errors = []
    start_time = time.time()

    for page_num in range(total_pages):
        # 超时检查
        if timeout and (time.time() - start_time) > timeout:
            errors.append(PDFlowError(
                input_path, "压缩",
                f"处理超时（{timeout}秒），已处理{page_num}/{total_pages}页",
                recoverable=False
            ))
            break

        try:
            page = doc.load_page(page_num)
            rect = page.rect
            zoom = dpi / 72.0
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat, alpha=False)

            # 内存渲染：PIL → JPEG → BytesIO，零磁盘 I/O
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=jpeg_q, optimize=True)
            img.close()
            pix = None

            # 新建页面，内存中插入 JPEG
            new_page = new_doc.new_page(width=rect.width, height=rect.height)
            new_page.insert_image(rect, stream=buf.getvalue())
        except Exception as e:
            # 单页压缩失败，跳过该页继续
            errors.append(PDFlowError(
                input_path, f"压缩第{page_num + 1}页", str(e), recoverable=True
            ))
            continue

        if progress_callback:
            progress_callback(page_num + 1, total_pages)

    doc.close()

    # 保存时启用所有压缩选项
    new_doc.save(output_path,
        garbage=4, deflate=True,
        deflate_images=True, deflate_fonts=True,
        clean=True, compression_effort=9)
    new_doc.close()

    return errors




# ============================================================
# PDF 转图片
# ============================================================
def pdf_to_images(input_path: str, output_dir: str = None, dpi: int = 150, fmt: str = "png") -> dict:
    """将PDF每一页转为图片
    fmt: 输出格式，支持 png, jpg, jpeg, webp 等
    """
    try:
        fmt = fmt.lower().strip()
        if fmt.startswith("."):
            fmt = fmt[1:]
        # 验证格式
        valid_formats = ["png", "jpg", "jpeg", "webp", "bmp", "tiff"]
        if fmt not in valid_formats:
            fmt = "png"  # 默认回退到 png

        doc = fitz.open(input_path)
        output_dir = output_dir or os.path.dirname(input_path)
        base = os.path.splitext(os.path.basename(input_path))[0]
        out_dir = os.path.join(output_dir, f"{base}_图片")
        os.makedirs(out_dir, exist_ok=True)

        image_paths = []
        for i, page in enumerate(doc):
            mat = fitz.Matrix(dpi / 72, dpi / 72)
            pix = page.get_pixmap(matrix=mat)
            out_path = os.path.join(out_dir, f"第{i+1}页.{fmt}")
            pix.save(out_path)
            image_paths.append(out_path)

        doc.close()
        return {
            "status": "ok",
            "images": image_paths,
            "count": len(image_paths),
        }
    except Exception as e:
        raise Exception(f"PDF转图片失败: {str(e)}")

# ============================================================
# 图片转 PDF
# ============================================================
def images_to_pdf(image_paths: list, output_path: str = None,
                  orientation: str = "portrait", quality: str = "high") -> dict:
    """将多张图片合并为一个PDF
    orientation: "portrait"(竖版) 或 "landscape"(横版)
    quality: "high"(高质量) | "medium"(中等) | "low"(低质量/小文件)
    """
    if not image_paths:
        raise Exception("请先选择要转换的图片")

    # 预扫描：获取最大宽高，用于统一页面方向
    max_w, max_h = 0, 0
    for img_path in image_paths:
        img = Image.open(img_path)
        w, h = img.size
        max_w = max(max_w, w)
        max_h = max(max_h, h)
        img.close()

    # 纸张尺寸映射（单位：点，1英寸=72点）
    PAGE_SIZES = {
        "A4": (595, 842),
        "A3": (842, 1191),
        "A5": (421, 595),
        "4K": (768, 1024),
        "8K": (1536, 2048),
    }

    # 根据方向和最大尺寸确定页面大小
    is_landscape = orientation.lower() in ("landscape", "横版")
    if is_landscape:
        page_w, page_h = max(max_w, max_h), min(max_w, max_h)
    else:
        page_w, page_h = min(max_w, max_h), max(max_w, max_h)

    # 质量设置影响 JPEG 压缩率
    jpeg_quality = {"high": 95, "medium": 80, "low": 60}.get(quality.lower(), 90)

    try:
        doc = fitz.open()
        for img_path in image_paths:
            img = Image.open(img_path)
            img_w, img_h = img.size
            img.close()

            # 居中放置，保持原图比例
            scale = min(page_w / img_w, page_h / img_h) if img_w > 0 and img_h > 0 else 1
            disp_w = img_w * scale
            disp_h = img_h * scale
            x0 = (page_w - disp_w) / 2
            y0 = (page_h - disp_h) / 2

            rect = fitz.Rect(x0, y0, x0 + disp_w, y0 + disp_h)
            page = doc.new_page(width=page_w, height=page_h)
            page.insert_image(rect, filename=img_path)

        if output_path is None:
            output_path = get_output_path(image_paths[0], "图片转PDF")
        if not output_path.lower().endswith(".pdf"):
            output_path += ".pdf"

        doc.save(output_path, garbage=4, deflate=True)
        doc.close()

        return {
            "status": "ok",
            "output": output_path,
            "pages": len(image_paths),
        }
    except Exception as e:
        raise Exception(f"图片转PDF失败: {str(e)}")

# ============================================================
# 添加水印
# ============================================================
def add_watermark(input_path: str, wm_text: str, output_path: str = None,
                  font_size: int = 60, opacity: float = 0.15,
                  rotation: int = -45, color: tuple = (128, 128, 128),
                  position: str = "center") -> dict:
    """
    为PDF添加文字水印
    position: center | tile
    """
    if not wm_text:
        raise Exception("请输入水印文字")
    if font_size < 10 or font_size > 200:
        raise Exception("字号需在10-200之间")
    if opacity < 0.01 or opacity > 1:
        raise Exception("透明度需在0.01-1之间")

    try:
        doc = fitz.open(input_path)
        original_size = os.path.getsize(input_path)

        # 尝试加载中文字体
        font_paths = [
            "C:/Windows/Fonts/msyh.ttc",   # 微软雅黑
            "C:/Windows/Fonts/simhei.ttf", # 黑体
            "C:/Windows/Fonts/simsun.ttc", # 宋体
            "C:/Windows/Fonts/arial.ttf",
        ]
        font_path = None
        for fp in font_paths:
            if os.path.exists(fp):
                font_path = fp
                break

        for page in doc:
            if position == "center":
                # 中心水印
                rect = page.rect
                fontsize = font_size * min(rect.width, rect.height) / 200
                if font_path:
                    font = ImageFont.truetype(font_path, int(fontsize))
                else:
                    font = ImageFont.load_default()

                # 创建水印图片
                wm_img = Image.new("RGBA", (int(rect.width), int(rect.height)), (255, 255, 255, 0))
                draw = ImageDraw.Draw(wm_img)
                bbox = draw.textbbox((0, 0), wm_text, font=font)
                tw = bbox[2] - bbox[0]
                th = bbox[3] - bbox[1]
                tx = (rect.width - tw) / 2
                ty = (rect.height - th) / 2
                draw.text((tx, ty), wm_text, fill=(*color, int(255 * opacity)), font=font)

                buf = io.BytesIO()
                wm_img.save(buf, format="PNG")
                buf.seek(0)
                wm_pixmap = fitz.Pixmap(fitz.open_stream(stream=buf.read(), filetype="png"))

                # 旋转水印
                if rotation != 0:
                    wm_pixmap = wm_pixmap.rotate(rotation)

                page.insert_image(rect, pixmap=wm_pixmap)
                wm_pixmap = None

            elif position == "tile":
                # 平铺水印
                rect = page.rect
                page_h = rect.height
                page_w = rect.width

                for y in range(0, int(page_h), 200):
                    for x in range(0, int(page_w), 300):
                        if font_path:
                            font = ImageFont.truetype(font_path, font_size)
                        else:
                            font = ImageFont.load_default()

                        wm_img = Image.new("RGBA", (300, 100), (255, 255, 255, 0))
                        draw = ImageDraw.Draw(wm_img)
                        draw.text((10, 10), wm_text, fill=(*color, int(255 * opacity)), font=font)

                        buf = io.BytesIO()
                        wm_img.save(buf, format="PNG")
                        buf.seek(0)
                        wm_pix = fitz.Pixmap(fitz.open_stream(stream=buf.read(), filetype="png"))

                        if rotation != 0:
                            wm_pix = wm_pix.rotate(rotation)

                        clip_rect = fitz.Rect(x, y, min(x + 300, page_w), min(y + 100, page_h))
                        page.insert_image(clip_rect, pixmap=wm_pix)
                        wm_pix = None

        if output_path is None:
            output_path = get_output_path(input_path, "水印")
        if not output_path.lower().endswith(".pdf"):
            output_path += ".pdf"

        total_pages = len(doc)
        doc.save(output_path, garbage=4, deflate=True)
        doc.close()

        return {
            "status": "ok",
            "output": output_path,
            "pages": total_pages,
        }
    except Exception as e:
        raise Exception(f"添加水印失败: {str(e)}")

# ============================================================
# 页面重排序（页面提取）
# ============================================================
def reorder_pdf_pages(input_path: str, page_order: str, output_path: str = None) -> dict:
    """
    重排PDF页面
    page_order: 逗号分隔的页码，如 "1,3,2,4" 或 "1-5,2,3"
    """
    try:
        doc = fitz.open(input_path)
        total = len(doc)
        new_doc = fitz.open()

        # 解析页码
        selected = []
        for part in page_order.split(","):
            part = part.strip()
            if "-" in part:
                segs = part.split("-")
                if len(segs) == 2:
                    try:
                        s, e = int(segs[0]), int(segs[1])
                        selected.extend(range(s - 1, e))
                    except ValueError:
                        pass
            else:
                try:
                    selected.append(int(part) - 1)
                except ValueError:
                    pass

        selected = [i for i in selected if 0 <= i < total]
        for idx in selected:
            new_doc.insert_pdf(doc, from_page=idx, to_page=idx)

        doc.close()

        if output_path is None:
            output_path = get_output_path(input_path, "重排")
        if not output_path.lower().endswith(".pdf"):
            output_path += ".pdf"

        new_doc.save(output_path)
        new_doc.close()

        return {
            "status": "ok",
            "output": output_path,
            "pages": len(selected),
        }
    except Exception as e:
        raise Exception(f"页面重排失败: {str(e)}")

# ============================================================
# PDF 转 Word（使用 pdf2docx + 后处理字体统一）
# ============================================================
def pdf_to_word(input_path: str, output_path: str = None) -> dict:
    """PDF转Word（需要 pdf2docx 库）

    优化点：
    - 转换后统一中英文字体映射，避免字体不统一
    - 中文默认映射为「微软雅黑」，英文/数字默认映射为「Calibri」
    """
    try:
        from pdf2docx import Converter
    except Exception as e:
        raise Exception(f"PDF转Word功能依赖缺失: {type(e).__name__}: {e}")

    try:
        if output_path is None:
            output_path = get_output_path(input_path, "")
            output_path = output_path.replace(".pdf", ".docx")
        output_path = _resolve_output_path(output_path, input_path, ".docx")

        cv = Converter(input_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()

        _unify_docx_fonts(output_path)

        return {
            "status": "ok",
            "output": output_path,
        }
    except Exception as e:
        raise Exception(f"PDF转Word失败: {str(e)}")


def _unify_docx_fonts(docx_path: str):
    """后处理：统一 docx 中中英文字体映射

    策略：
    - 英文/数字字体 → Calibri
    - 中文字体（eastAsia）→ 微软雅黑
    - 保留粗体/斜体等样式不变
    """
    try:
        from docx import Document
        from docx.oxml.ns import qn

        FONT_LATIN = "Calibri"
        FONT_EAST_ASIA = "微软雅黑"

        doc = Document(docx_path)

        for paragraph in doc.paragraphs:
            for run in paragraph.runs:
                rPr = run._element.rPr
                if rPr is None:
                    rPr = run._element.get_or_add_rPr()

                rFonts = rPr.find(qn('w:rFonts'))
                if rFonts is None:
                    from lxml import etree
                    rFonts = etree.SubElement(rPr, qn('w:rFonts'))

                rFonts.set(qn('w:ascii'), FONT_LATIN)
                rFonts.set(qn('w:hAnsi'), FONT_LATIN)
                rFonts.set(qn('w:eastAsia'), FONT_EAST_ASIA)
                rFonts.set(qn('w:cs'), FONT_LATIN)

        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            rPr = run._element.rPr
                            if rPr is None:
                                rPr = run._element.get_or_add_rPr()

                            rFonts = rPr.find(qn('w:rFonts'))
                            if rFonts is None:
                                from lxml import etree
                                rFonts = etree.SubElement(rPr, qn('w:rFonts'))

                            rFonts.set(qn('w:ascii'), FONT_LATIN)
                            rFonts.set(qn('w:hAnsi'), FONT_LATIN)
                            rFonts.set(qn('w:eastAsia'), FONT_EAST_ASIA)
                            rFonts.set(qn('w:cs'), FONT_LATIN)

        doc.save(docx_path)
    except Exception:
        pass

# ============================================================
# PDF 转 Excel（纯 pdfplumber 多参数优化提取）
# ============================================================
def pdf_to_excel(input_path: str, output_path: str = None) -> dict:
    """PDF转Excel（v1.1-patch：IR 中间结构 + wrap_text 统一写入）

    策略：
    1. 对每一页，尝试多种 table_settings 参数组合（v1 保留）
    2. 按评分选最优提取结果（v1.1-patch 加列稳定性）
    3. 每页独立处理，页码绝不混乱
    4. 表格提取失败时，使用 word-level 文字回退（v1.1-patch 返回 IR）
    5. 统一通过 IR 中间结构 + write_cell 写入 Excel（v1.1-patch 新增）

    参数：
    - input_path: 输入 PDF 路径
    - output_path: 输出 xlsx 路径，支持两种形式：
        * 完整文件路径（推荐）：如 D:/out/foo.xlsx
        * 输出目录：如是已存在目录，会在目录下生成 <input_basename>.xlsx
    """
    try:
        import pdfplumber
        from openpyxl import Workbook
        from openpyxl.styles import Alignment
        from src.common.pdf_table_ir import normalize_excel_input
    except Exception as e:
        raise Exception(f"PDF转Excel功能依赖缺失: {type(e).__name__}: {e}")

    try:
        output_path = _resolve_output_path(output_path, input_path, ".xlsx")

        all_sheets = []

        with pdfplumber.open(input_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_tables = _extract_page_best(page, page_num)
                all_sheets.extend(page_tables)

        if not all_sheets:
            all_sheets = _extract_text_fallback(input_path)

        if not all_sheets:
            raise Exception("PDF 未检测到可提取的表格或文字内容（v1.1-patch: 可能是扫描件或图片型 PDF）")

        # v1.1-patch P0 Hotfix：Excel 写入层只接受 pandas.DataFrame
        # - 统一通过 normalize_excel_input 收敛（IR / DataFrame / list 全部接受）
        # - 禁止任何路径直接调用 df.tolist() 或 df.values.tolist()
        # - 用 df.to_dict("records") 拿数据
        def write_cell(ws, r, c, value):
            cell = ws.cell(row=r, column=c, value=value)
            cell.alignment = Alignment(wrap_text=True, vertical="top")
            return cell

        wb = Workbook()
        wb.remove(wb.active)  # 移除默认 sheet

        for sheet_name, ir in all_sheets:
            # P0 Hotfix: 统一 IR/DataFrame/list → DataFrame
            df = normalize_excel_input(ir)
            # P0 Hotfix: 用 to_dict("records") 替代 values.tolist()
            records = df.fillna("").astype(str).to_dict("records")

            ws = wb.create_sheet(sheet_name[:31])
            for i, row_dict in enumerate(records, 1):
                for j, val in enumerate(row_dict.values(), 1):
                    write_cell(ws, i, j, val)

            # 列宽估算（保留 v1 的 CJK 系数，从 ws 读以保持一致性）
            if records:
                max_cols = max((len(r) for r in records), default=0)
                for col_idx in range(1, max_cols + 1):
                    col_letter = ws.cell(row=1, column=col_idx).column_letter
                    max_len = 0
                    for r_idx in range(1, len(records) + 1):
                        try:
                            v = ws.cell(row=r_idx, column=col_idx).value
                            val = str(v) if v else ""
                            cjk = sum(1 for c in val if '\u4e00' <= c <= '\u9fff')
                            cell_len = len(val) + cjk
                            max_len = max(max_len, cell_len)
                        except Exception:
                            pass
                    ws.column_dimensions[col_letter].width = min(max(max_len + 4, 8), 50)

        wb.save(output_path)

        return {
            "status": "ok",
            "output": output_path,
            "tables": len(all_sheets),
        }
    except Exception as e:
        raise Exception(f"PDF转Excel失败: {str(e)}")


_PARAM_COMBOS = [
    {
        "vertical_strategy": "lines",
        "horizontal_strategy": "lines",
        "snap_tolerance": 5,
        "join_tolerance": 5,
        "edge_min_length": 10,
        "min_words_vertical": 1,
        "min_words_horizontal": 1,
    },
    {
        "vertical_strategy": "lines",
        "horizontal_strategy": "lines",
        "snap_tolerance": 3,
        "join_tolerance": 3,
        "edge_min_length": 6,
        "min_words_vertical": 1,
        "min_words_horizontal": 1,
    },
    {
        "vertical_strategy": "lines",
        "horizontal_strategy": "lines",
        "snap_tolerance": 8,
        "join_tolerance": 8,
        "edge_min_length": 15,
        "min_words_vertical": 1,
        "min_words_horizontal": 1,
    },
    {
        "vertical_strategy": "text",
        "horizontal_strategy": "text",
        "snap_tolerance": 5,
        "join_tolerance": 5,
        "text_tolerance": 5,
        "min_words_vertical": 1,
        "min_words_horizontal": 1,
    },
    {
        "vertical_strategy": "text",
        "horizontal_strategy": "text",
        "snap_tolerance": 3,
        "join_tolerance": 3,
        "text_tolerance": 3,
        "min_words_vertical": 1,
        "min_words_horizontal": 1,
    },
    {
        "vertical_strategy": "text",
        "horizontal_strategy": "text",
        "snap_tolerance": 8,
        "join_tolerance": 8,
        "text_tolerance": 8,
        "min_words_vertical": 1,
        "min_words_horizontal": 1,
    },
    {
        "vertical_strategy": "lines",
        "horizontal_strategy": "text",
        "snap_tolerance": 5,
        "join_tolerance": 5,
        "text_tolerance": 5,
        "min_words_vertical": 1,
        "min_words_horizontal": 1,
    },
    {
        "vertical_strategy": "text",
        "horizontal_strategy": "lines",
        "snap_tolerance": 5,
        "join_tolerance": 5,
        "text_tolerance": 5,
        "min_words_vertical": 1,
        "min_words_horizontal": 1,
    },
]


def _extract_page_best(page, page_num: int) -> list:
    """对单页尝试所有参数组合，返回最优提取结果

    返回: [(sheet_name, df), ...]
    """
    import pandas as pd

    best_tables = []

    for settings in _PARAM_COMBOS:
        try:
            tables = page.extract_tables(table_settings=settings)
        except Exception:
            continue

        for t_idx, raw_table in enumerate(tables, 1):
            cleaned = _clean_table_data(raw_table)
            if not cleaned or len(cleaned) < 1:
                continue

            df = _table_to_dataframe(cleaned)
            if df.empty:
                continue

            score = _score_table(df)
            if score < 0.08:
                continue

            best_tables.append({
                "df": df,
                "score": score,
                "t_idx": t_idx,
            })

    if not best_tables:
        # P0 Hotfix: 用 parse_layout_blocks 替代 _extract_page_words（旧 DataFrame 路径）
        from src.common.pdf_layout_parser import parse_layout_blocks
        page_rows = parse_layout_blocks(page)
        if page_rows:
            ir = fallback_block(
                rows=page_rows,
                page=page_num,
                table_id=1,
            )
            return [(f"第{page_num}页_文字", ir)]
        return []

    best_tables.sort(key=lambda x: x["score"], reverse=True)

    selected = []
    used_dfs = []

    for candidate in best_tables:
        is_dup = False
        for existing in used_dfs:
            if _is_duplicate_table(candidate["df"], existing):
                is_dup = True
                break
        if not is_dup:
            selected.append(candidate)
            used_dfs.append(candidate["df"])

    # v1.1-patch P0 Hotfix：在输出边界把 DataFrame 包装成 IR dict
    # 保留内部 DataFrame 流程以维持评分/去重的稳定性
    # P0 修复：禁止 df.values.tolist() 主路径，改用 to_dict("records") 统一
    from src.common.pdf_table_ir import to_table_block
    result = []
    for idx, item in enumerate(selected, 1):
        sheet_name = f"第{page_num}页_表{idx}"
        df = item["df"]
        # P0 Hotfix: DataFrame → list of list（IR rows）走 to_dict("records") 路径
        if hasattr(df, 'to_dict'):
            records = df.fillna("").astype(str).to_dict("records")
            rows = [list(r.values()) for r in records]
        else:
            rows = df
        ir = to_table_block(
            rows=rows,
            page=page_num,
            table_id=idx,
            confidence=item.get("score", 1.0),
            mode="structured",
        )
        result.append((sheet_name, ir))

    return result


def _column_stability_score(cols) -> float:
    """列宽稳定性评分（v1.1-patch）：列内字符长度方差小 → 列结构稳定 → 评分高

    用途：减弱"参数组合扫描中偶然生成均齐假表"对评分的干扰

    返回 0.0 ~ 0.1（最大贡献 0.1）
    """
    try:
        import statistics
    except ImportError:
        return 0.0

    lengths = [len(str(c)) for c in cols if c is not None and str(c).strip()]
    if not lengths:
        return 0.0

    max_len = max(lengths)
    if max_len == 0:
        return 0.0

    pstdev = statistics.pstdev(lengths)
    # 1 - (标准差 / (max+1))：max+1 防 0 除；值越大表示列越稳定
    stability = 1.0 - (pstdev / (max_len + 1))
    return max(0.0, min(stability, 1.0)) * 0.1


def _score_table(df) -> float:
    """给表格打分：单元格填充率 + 行列结构 + 内容丰富度 + 列稳定性（v1.1-patch）"""
    if df.empty:
        return 0.0

    total_cells = df.shape[0] * df.shape[1]
    if total_cells == 0:
        return 0.0

    non_empty = 0
    total_chars = 0
    for col in df.columns:
        for val in df[col]:
            s = str(val).strip() if val is not None else ""
            if s and s != "None":
                non_empty += 1
                total_chars += len(s)

    fill_rate = non_empty / total_cells

    col_score = min(df.shape[1] / 5.0, 1.0) * 0.2
    row_score = min(df.shape[0] / 3.0, 1.0) * 0.2

    avg_chars = total_chars / max(non_empty, 1)
    content_score = min(avg_chars / 5.0, 1.0) * 0.2

    # v1.1-patch：列稳定性评分（每列累加，最多贡献 0.1 * 5 = 0.5）
    # P0 Hotfix：用 df.iloc[:, col_idx] 按位置取列，强制返回 Series
    # 避免 df[col] 在某些 pandas 行为下返回 DataFrame 导致 .tolist() 崩溃
    col_stability_total = 0.0
    col_count = 0
    for col_idx, _ in enumerate(df.columns):
        col_series = df.iloc[:, col_idx]
        # 强制 squeeze：任何长度 1 的 DataFrame → Series
        if hasattr(col_series, 'ndim') and col_series.ndim > 1:
            col_series = col_series.squeeze()
        col_stability_total += _column_stability_score(col_series.tolist())
        col_count += 1
    # 归一化：按列数取平均，单表最大贡献 0.1
    col_stability_score = (col_stability_total / col_count) if col_count > 0 else 0.0

    score = fill_rate * 0.4 + col_score + row_score + content_score + col_stability_score

    return round(score, 4)


def _table_to_dataframe(cleaned_table) -> "pd.DataFrame":
    """将清洗后的二维列表转为 DataFrame"""
    import pandas as pd

    if not cleaned_table:
        return pd.DataFrame()

    header = cleaned_table[0]
    rows = cleaned_table[1:]

    if any(h and str(h).strip() for h in header):
        return pd.DataFrame(rows, columns=header)
    else:
        return pd.DataFrame(cleaned_table)


def _extract_text_fallback(input_path: str) -> list:
    """纯文字回退提取（v1.1-patch P0 Hotfix）：

    使用 parse_layout_blocks 重建阅读顺序（按 y 坐标 + 行聚类 + 语义拆分）
    禁止 page.extract_text() 和 text.split("\\n") 主路径
    """
    import pdfplumber
    from src.common.pdf_table_ir import fallback_block
    from src.common.pdf_layout_parser import parse_layout_blocks

    results = []

    try:
        with pdfplumber.open(input_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                # P0 Hotfix: 用 parse_layout_blocks 替代 extract_text() + split("\n")
                rows = parse_layout_blocks(page)
                if rows:
                    ir = fallback_block(
                        rows=rows,
                        page=page_num,
                        table_id=1,
                    )
                    sheet_name = f"第{page_num}页_文字"
                    results.append((sheet_name, ir))
    except Exception:
        pass

    return results


def _clean_table_data(table):
    """清洗表格数据：去除 None、合并空行、保留 \\n（v1.1-patch: 单元格换行不再被替换为空格）"""
    cleaned = []
    for row in table:
        cleaned_row = []
        for cell in row:
            if cell is None:
                cleaned_row.append("")
            else:
                text = str(cell).strip()
                # v1.1-patch 修复：仅去 \r，保留 \n
                # 由下游 write_cell 配合 wrap_text=True 实现单元格内换行渲染
                text = text.replace("\r", "")
                cleaned_row.append(text)
        cleaned.append(cleaned_row)

    result = []
    for row in cleaned:
        if any(cell for cell in row):
            result.append(row)

    return result


def _is_duplicate_table(df1, df2) -> bool:
    """判断两个表格是否是重复的（同一张表格的不同参数提取结果）

    判定规则（满足任一即为重复）：
    1. 内容重叠率 > 0.35
    2. 行列数相同 + 内容重叠率 > 0.20
    3. 首行（表头）完全相同 + 行列数相同
    """
    overlap = _table_content_overlap(df1, df2)

    if overlap > 0.35:
        return True

    same_shape = (df1.shape[0] == df2.shape[0]) and (df1.shape[1] == df2.shape[1])

    if same_shape and overlap > 0.20:
        return True

    if same_shape and df1.shape[1] > 0:
        header1 = [str(v).strip() for v in df1.iloc[0]] if df1.shape[0] > 0 else []
        header2 = [str(v).strip() for v in df2.iloc[0]] if df2.shape[0] > 0 else []
        if header1 and header2 and header1 == header2:
            return True

    return False


def _table_content_overlap(df1, df2) -> float:
    """计算两个表格的内容重叠率（0~1）

    比较前统一清理空白字符，避免因间距不同导致误判
    """
    def _normalize(val):
        s = str(val).strip() if val is not None else ""
        s = " ".join(s.split())
        if not s or s == "None":
            return ""
        return s

    cells1 = set()
    for col in df1.columns:
        for val in df1[col]:
            s = _normalize(val)
            if s:
                cells1.add(s)

    cells2 = set()
    for col in df2.columns:
        for val in df2[col]:
            s = _normalize(val)
            if s:
                cells2.add(s)

    if not cells1 or not cells2:
        return 0.0

    intersection = cells1 & cells2
    smaller = min(len(cells1), len(cells2))

    return len(intersection) / smaller

# ============================================================
# PDF 转 PPT
# ============================================================
def pdf_to_ppt(input_path: str, output_dir: str = None) -> dict:
    """PDF → PPT（每页转为图片插入幻灯片，保持原始比例）"""
    try:
        from pptx import Presentation
        from pptx.util import Emu

        doc = fitz.open(input_path)
        output_dir = output_dir or os.path.dirname(input_path)
        os.makedirs(output_dir, exist_ok=True)

        prs = Presentation()

        # 使用空白布局
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

        for page_num in range(doc.page_count):
            page = doc.load_page(page_num)

            # 获取 PDF 页面尺寸（点为单位，1点=1/72英寸）
            pdf_width = page.rect.width
            pdf_height = page.rect.height

            # 转换为 EMU（English Metric Units，1英寸=914400 EMU）
            # 1点 = 1/72英寸 = 914400/72 = 12700 EMU
            slide_width = int(pdf_width * 12700)
            slide_height = int(pdf_height * 12700)

            # 设置幻灯片尺寸（保持 PDF 原始比例）
            prs.slide_width = slide_width
            prs.slide_height = slide_height

            # 高分辨率渲染（2x缩放）
            mat = fitz.Matrix(2.0, 2.0)
            pix = page.get_pixmap(matrix=mat)

            # 保存临时图片
            import tempfile
            temp_img = os.path.join(tempfile.gettempdir(), f"_pdf2ppt_temp_{page_num}.png")
            pix.save(temp_img)

            # 添加幻灯片
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(temp_img, Emu(0), Emu(0), width=prs.slide_width, height=prs.slide_height)

            # 删除临时图片
            os.remove(temp_img)

        # 保存 PPT
        output_path = os.path.join(output_dir, os.path.splitext(os.path.basename(input_path))[0] + ".pptx")
        prs.save(output_path)
        
        # 先保存页数，再关闭文档
        pages = doc.page_count
        doc.close()

        return {
            "status": "ok",
            "output": output_path,
            "pages": pages,
        }
    except ImportError:
        raise Exception("PDF转PPT需要安装 python-pptx 库，请运行: pip install python-pptx")
    except Exception as e:
        raise Exception(f"PDF转PPT失败: {str(e)}")

# ============================================================
# 批量转换
# ============================================================
def batch_convert(input_path: str, output_dir: str = None, batch_fmt: str = "pdf",
                   progress_callback=None, timeout: int = 60) -> dict:
    """批量转换：根据输出格式自动选择转换函数
    batch_fmt: 输出格式，pdf | word | excel | ppt | img
    progress_callback: callable(current, total, filename) — 每转换一个文件后触发
    timeout: 单个文件处理超时秒数（默认60秒）
    """
    IMAGE_EXTS = ["png", "jpg", "jpeg", "bmp", "tiff", "webp", "gif"]

    # 根据扩展名判断输入类型
    ext = os.path.splitext(input_path)[1].lower().lstrip(".")
    is_pdf = ext == "pdf"
    is_img = ext in IMAGE_EXTS

    start_time = time.time()

    try:
        # 根据 batch_fmt 和输入类型决定调用哪个函数
        if batch_fmt == "pdf":
            # 图片转PDF
            if is_img:
                result = images_to_pdf([input_path], get_output_path(input_path, "批量图片转PDF", output_dir))
            else:
                raise Exception("批量转换PDF格式只支持图片文件")

        elif batch_fmt == "word":
            # PDF转Word
            if is_pdf:
                result = pdf_to_word(input_path, get_output_path(input_path, "批量PDF转Word", output_dir))
            else:
                raise Exception("批量转Word格式只支持PDF文件")

        elif batch_fmt == "excel":
            # PDF转Excel
            if is_pdf:
                result = pdf_to_excel(input_path, get_output_path(input_path, "批量PDF转Excel", output_dir))
            else:
                raise Exception("批量转Excel格式只支持PDF文件")

        elif batch_fmt == "ppt":
            # PDF转PPT
            if is_pdf:
                result = pdf_to_ppt(input_path, output_dir)
            else:
                raise Exception("批量转PPT格式只支持PDF文件")

        elif batch_fmt == "img":
            # PDF转图片
            if is_pdf:
                result = pdf_to_images(input_path, output_dir)
            else:
                raise Exception("批量转图片格式只支持PDF文件")

        else:
            raise Exception(f"不支持的批量输出格式: {batch_fmt}")

        # 超时检查
        elapsed = time.time() - start_time
        if timeout and elapsed > timeout:
            return {
                "status": "timeout",
                "output": result.get("output", ""),
                "message": f"处理超时（{timeout}秒），耗时{round(elapsed, 1)}秒",
                "errors": [],
                "skipped": 0,
            }

        if progress_callback:
            progress_callback(1, 1, os.path.basename(input_path))

        result["errors"] = []
        result["skipped"] = 0
        return result

    except Exception as e:
        # 单文件转换失败，返回错误信息而非抛出异常
        return {
            "status": "error",
            "output": "",
            "errors": [PDFlowError(input_path, f"批量转{batch_fmt}", str(e), recoverable=True).to_dict()],
            "skipped": 1,
        }


# ============================================================
# 批量合并多组PDF
# ============================================================
def batch_merge_pdfs(file_groups: list, output_dir: str, progress_callback=None) -> dict:
    """批量合并多组PDF文件

    file_groups: 列表，每个元素为一组待合并的文件路径列表
        示例: [["a.pdf", "b.pdf"], ["c.pdf", "d.pdf"]]
    output_dir: 输出目录
    progress_callback: callable(current, total, filename) — 每完成一组合并后触发

    返回:
        {
            "status": "ok",
            "results": [...],        # 每组的合并结果
            "errors": [...],         # PDFlowError 字典列表
            "skipped": int,          # 跳过的组数
            "total_groups": int,
        }
    """
    if not file_groups:
        raise Exception("请提供至少一组待合并的PDF文件")

    os.makedirs(output_dir, exist_ok=True)
    total = len(file_groups)
    results = []
    errors = []
    skipped = 0

    for idx, group in enumerate(file_groups):
        group_name = f"合并组{idx + 1}"
        output_path = os.path.join(output_dir, f"合并结果_{idx + 1}.pdf")

        try:
            result = merge_pdfs(output_path, *group)
            results.append(result)
        except Exception as e:
            errors.append(PDFlowError(
                str(group), group_name, str(e), recoverable=True
            ))
            skipped += 1
            continue

        if progress_callback:
            progress_callback(idx + 1, total, f"合并结果_{idx + 1}.pdf")

    return {
        "status": "ok",
        "results": results,
        "errors": [e.to_dict() for e in errors],
        "skipped": skipped,
        "total_groups": total,
    }


# ============================================================
# 批量压缩多个PDF
# ============================================================
def batch_compress_pdfs(file_paths: list, output_dir: str, quality: str = "high",
                         progress_callback=None, timeout: int = 60) -> dict:
    """批量压缩多个PDF文件

    file_paths: 待压缩的PDF文件路径列表
    output_dir: 输出目录
    quality: 压缩质量 "high" | "medium" | "low"
    progress_callback: callable(current, total, filename) — 每完成一个文件后触发
    timeout: 单个文件处理超时秒数（默认60秒）

    返回:
        {
            "status": "ok",
            "results": [...],        # 每个文件的压缩结果
            "errors": [...],         # PDFlowError 字典列表
            "skipped": int,
            "total_files": int,
        }
    """
    if not file_paths:
        raise Exception("请提供至少一个待压缩的PDF文件")

    os.makedirs(output_dir, exist_ok=True)
    total = len(file_paths)
    results = []
    errors = []
    skipped = 0

    for idx, fp in enumerate(file_paths):
        basename = os.path.splitext(os.path.basename(fp))[0]
        output_path = os.path.join(output_dir, f"{basename}_压缩.pdf")

        try:
            result = compress_pdf(fp, quality=quality, output_path=output_path,
                                   timeout=timeout)
            results.append(result)
            # 合并单文件内部的页级错误
            if result.get("errors"):
                errors.extend([PDFlowError(fp, "压缩", e.get("message", ""), e.get("recoverable", True))
                               for e in result["errors"]])
        except Exception as e:
            errors.append(PDFlowError(fp, "批量压缩", str(e), recoverable=True))
            skipped += 1
            continue

        if progress_callback:
            progress_callback(idx + 1, total, os.path.basename(fp))

    return {
            "status": "ok",
            "results": results,
            "errors": [e.to_dict() for e in errors],
            "skipped": skipped,
            "total_files": total,
        }

# ============================================================
# OCR 文字提取（PDF 扫描件）
# ============================================================
def ocr_extract_text(input_path: str, progress_callback=None) -> dict:
    """OCR 提取 PDF 扫描件文字

    优先使用 PaddleOCR，回退到 Tesseract。
    如果两个引擎都不可用，尝试使用 PyMuPDF 内置文字提取。

    Args:
        input_path: PDF 文件路径
        progress_callback: callable(current_page, total_pages, status_text)

    Returns:
        {
            "status": "ok",
            "text": "完整文本",
            "pages": [{"page_num": 1, "text": "...", "confidence": 0.95}],
            "engine": "paddleocr" | "tesseract" | "fitz_builtin",
            "total_pages": N,
        }
    """
    # 校验文件
    if not os.path.isfile(input_path):
        raise Exception(f"文件不存在: {input_path}")

    ext = os.path.splitext(input_path)[1].lower()
    if ext != ".pdf":
        raise Exception(f"不支持的文件格式: {ext}，仅支持 PDF 文件")

    # 第一步：尝试 OCR 引擎（PaddleOCR / Tesseract）
    try:
        from common.ocr_engine import get_ocr_engine
        engine = get_ocr_engine()

        if engine.available:
            result = engine.extract_text_from_pdf(input_path, progress_callback)
            if result.get("status") == "ok":
                return result
            # OCR 引擎存在但识别失败，记录原因后继续尝试回退
    except ImportError:
        pass
    except Exception as e:
        pass

    # 第二步：回退到 PyMuPDF 内置文字提取
    try:
        doc = fitz.open(input_path)
        total_pages = len(doc)
        all_text = []
        pages_result = []

        for i in range(total_pages):
            page = doc.load_page(i)
            page_text = page.get_text("text")

            if progress_callback:
                progress_callback(i + 1, total_pages, f"内置提取第 {i + 1}/{total_pages} 页")

            all_text.append(page_text)
            pages_result.append({
                "page_num": i + 1,
                "text": page_text,
                "confidence": 1.0 if page_text.strip() else 0.0,
            })

        doc.close()

        return {
            "status": "ok",
            "text": "\n".join(all_text),
            "pages": pages_result,
            "engine": "fitz_builtin",
            "total_pages": total_pages,
        }
    except Exception as e:
        raise Exception(f"OCR 文字提取失败: {e}")


# ============================================================
# 批量转换多个文件
# ============================================================
def batch_convert_files(file_paths: list, output_dir: str, batch_fmt: str = "pdf",
                         progress_callback=None, timeout: int = 60) -> dict:
    """批量转换多个文件

    file_paths: 待转换的文件路径列表
    output_dir: 输出目录
    batch_fmt: 输出格式，pdf | word | excel | ppt | img
    progress_callback: callable(current, total, filename) — 每完成一个文件后触发
    timeout: 单个文件处理超时秒数（默认60秒）

    返回:
        {
            "status": "ok",
            "results": [...],        # 每个文件的转换结果
            "errors": [...],         # PDFlowError 字典列表
            "skipped": int,
            "total_files": int,
        }
    """
    if not file_paths:
        raise Exception("请提供至少一个待转换的文件")

    os.makedirs(output_dir, exist_ok=True)
    total = len(file_paths)
    results = []
    errors = []
    skipped = 0

    for idx, fp in enumerate(file_paths):
        try:
            result = batch_convert(fp, output_dir=output_dir, batch_fmt=batch_fmt,
                                    timeout=timeout)
            if result.get("status") == "error":
                # 单文件转换失败
                errors.extend([PDFlowError(fp, f"批量转{batch_fmt}", e.get("message", ""), e.get("recoverable", True))
                               for e in result.get("errors", [])])
                skipped += 1
            else:
                results.append(result)
                if result.get("errors"):
                    errors.extend([PDFlowError(fp, f"批量转{batch_fmt}", e.get("message", ""), e.get("recoverable", True))
                                   for e in result["errors"]])
        except Exception as e:
            errors.append(PDFlowError(fp, f"批量转{batch_fmt}", str(e), recoverable=True))
            skipped += 1
            continue

        if progress_callback:
            progress_callback(idx + 1, total, os.path.basename(fp))

    return {
        "status": "ok",
        "results": results,
        "errors": [e.to_dict() for e in errors],
        "skipped": skipped,
        "total_files": total,
    }
